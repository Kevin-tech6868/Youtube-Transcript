import streamlit as st
from pytube import Search
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from pptx.util import Inches
from gtts import gTTS
from transformers import pipeline
import os
import tempfile

# Helper functions
def get_youtube_links(query):
    """Fetch YouTube links related to the query."""
    search = Search(query)
    videos = search.results
    return [{"title": video.title, "url": video.watch_url} for video in videos[:5]]  # Return top 5 links

def extract_transcript(video_url):
    """Extract transcript from a YouTube video."""
    video_id = video_url.split('watch?v=')[-1]
    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript = transcript_list.find_transcript(['en'])  # Get English transcript
        return " ".join([entry['text'] for entry in transcript.fetch()])
    except Exception as e:
        print(f"Error: {e}")
        return None  # Return None if there was an error

def summarize_text(text):
    """Summarize the given text using Hugging Face transformers."""
    summarizer = pipeline('summarization')
    summarized = summarizer(text, max_length=150, min_length=30, do_sample=False)
    return summarized[0]['summary_text']

def text_to_audio(text, output_path):
    """Convert text to speech and save as audio file."""
    tts = gTTS(text=text, lang='en')
    tts.save(output_path)

def create_ppt_from_text(text, output_ppt):
    """Create a PowerPoint presentation from text."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    
    # Split the text into slides (for simplicity, splitting by sentences)
    for sentence in text.split(". "):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]
        
        title.text = "Generated Content"
        content.text = sentence
    
    prs.save(output_ppt)

def convert_ppt_to_video(ppt_path, output_video):
    """Convert PowerPoint to video using moviepy."""
    clip = VideoFileClip(ppt_path)
    clip.write_videofile(output_video)

# Streamlit app
st.title("YouTube Summarizer with video")

# Step 1: User Input
user_input = st.text_input("Enter your query for YouTube video search:", "")

if user_input:
    # Step 2: Fetch YouTube links
    st.write("Fetching YouTube links...")
    video_links = get_youtube_links(user_input)
    
    for i, video in enumerate(video_links):
        st.write(f"{i+1}. [{video['title']}]({video['url']})")
    
    selected_link = st.text_input("Paste the URL of the video you want to use:")

    if selected_link:
        # Step 3: Extract video transcript
        st.write("Extracting transcript...")
        transcript = extract_transcript(selected_link)
        
        if transcript:
            st.write("Transcript extracted successfully.")
            st.text_area("Generated Transcript", transcript, height=200)
            
            # Step 4: Summarize the transcript
            st.write("Summarizing the transcript...")
            summarized_text = summarize_text(transcript)
            st.write("Summary generated successfully.")
            st.text_area("Summarized Content", summarized_text, height=200)

            # Step 5: Convert summarized transcript to audio
            audio_file = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3").name
            st.write("Converting summarized transcript to audio...")
            text_to_audio(summarized_text, audio_file)
            st.audio(audio_file)

            # Step 6: Create PowerPoint from summarized transcript
            ppt_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
            st.write("Creating PowerPoint from summarized transcript...")
            create_ppt_from_text(summarized_text, ppt_file)
            st.write("PowerPoint created. Download it below.")
            with open(ppt_file, "rb") as ppt:
                st.download_button("Download PPT", ppt, file_name="presentation.pptx")

            # Step 7: Convert PowerPoint to Video
            video_file = tempfile.NamedTemporaryFile(delete=False, suffix=".mp4").name
            st.write("Converting PowerPoint to video...")
            convert_ppt_to_video(ppt_file, video_file)
            st.write("Video created. Download it below.")
            with open(video_file, "rb") as video:
                st.download_button("Download Video", video, file_name="presentation.mp4")
        else:
            st.error("Unable to extract transcript. The video may not have subtitles or an error occurred.")
